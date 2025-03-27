import re
import os
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
import gptText
import addphoto

# Define template paths
TEMPLATES = {
    1: "template/minimalistic.pptx",  # Replace with the actual path to your first template
    2: "template/colourful.pptx",  # Replace with the actual path to your second template
    3: "template/professional.pptx",  # Replace with the actual path to your third template
    4: "template/dark.pptx",  # Replace with the actual path to your fourth template
}

def gettext(topicList, code: bool):
    slides_data = gptText.structured(topic_list=topicList, include_code=code)
    if slides_data and isinstance(slides_data, list) and len(slides_data) > 0 and 'Slides' in slides_data[0]:
        return slides_data[0]['Slides']
    else:
        return ""

def getphoto(SlideData):
    pattern = r"Image Suggestion:\s*(.+)"
    image_suggestions = re.findall(pattern, SlideData)
    image_paths = []
    for suggestion in image_suggestions:
        try:
            image_result = addphoto.get_images(suggestion, 1)
            if image_result:
                image_paths.append(image_result[0])
            else:
                print(f"No image found for suggestion: {suggestion}")
        except Exception as e:
            print(f"Error getting image for suggestion '{suggestion}': {e}")
    return image_paths

def create_presentation(slide_data, image_paths=None, output_filename="presentation.pptx", template_path=None):
    if image_paths is None:
        image_paths = []

    # Load the selected template or create a blank presentation
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        # Clear all existing slides
        xml_slides = prs.slides._sldIdLst  # Access the slide ID list
        slides = list(xml_slides)
        for slide in slides:
            xml_slides.remove(slide)  # Remove each slide
    else:
        prs = Presentation()

    # Add a blank slide layout to start fresh
    blank_slide_layout = prs.slide_layouts[6]

    # Extract the title from the slide data
    title_match = re.search(r"Title:\s*(.*)", slide_data)
    if title_match:
        overall_title = title_match.group(1).strip()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        # Check if the title placeholder exists
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = overall_title
        else:
            # Add a new text box for the title if the placeholder doesn't exist
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = overall_title
            title_p = title_frame.paragraphs[0]
            title_p.font.bold = True
            title_p.font.size = Pt(32)
            title_p.alignment = PP_ALIGN.CENTER

        # Check if a subtitle placeholder exists
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = "Generated Presentation"

    # Split the slide data into individual slides
    slides = re.split(r"---", slide_data)
    image_index = 0

    for slide_content in slides:
        slide_content = slide_content.strip()
        if not slide_content:
            continue

        # Extract the slide title
        title_match = re.search(r"Slide\s*\d+:\s*(.*)", slide_content)
        slide_title = title_match.group(1).strip() if title_match else "Slide"

        # Add a new slide using the blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add the title to the slide
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = slide_title
        title_p = title_frame.paragraphs[0]
        title_p.font.bold = True
        title_p.font.size = Pt(32)
        title_p.alignment = PP_ALIGN.CENTER

        # Add the content to the slide
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(6), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        lines = slide_content.split("\n")
        for line in lines:
            line = line.strip()

            if line.startswith("Slide") and ":" in line:
                continue
            elif line.startswith("Image Suggestion:"):
                continue
            elif line.startswith("-"):
                p = tf.add_paragraph()
                p.text = line[1:].strip()
                p.level = 0
                p.font.size = Pt(20)
            elif line.startswith("```"):
                code_match = re.search(r"```[a-zA-Z]*\n(.*?)```", slide_content, re.DOTALL)
                if code_match:
                    code = code_match.group(1).strip()
                    code_slide = prs.slides.add_slide(blank_slide_layout)
                    code_title_box = code_slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
                    code_title_frame = code_title_box.text_frame
                    code_title_frame.text = "Code Example"
                    code_title_p = code_title_frame.paragraphs[0]
                    code_title_p.font.bold = True
                    code_title_p.font.size = Pt(32)
                    code_title_p.alignment = PP_ALIGN.CENTER

                    code_box = code_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
                    code_tf = code_box.text_frame
                    code_tf.word_wrap = True
                    code_tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    p = code_tf.add_paragraph()
                    p.text = code
                    p.font.name = 'Courier New'
                    p.font.size = Pt(10)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.fill.solid()
                    p.font.fill.fore_color.rgb = RGBColor(40, 40, 40)
            elif line:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(20)

        # Add an image to the slide if available
        if image_paths and image_index < len(image_paths):
            try:
                img_path = image_paths[image_index]

                # Define the dimensions for the image
                left = Inches(6.5)  # Place the image on the right side
                top = Inches(1)  # Align the image with the top of the content
                width = Inches(3)  # Elongated width
                height = Inches(5)  # Elongated height

                # Add the image to the slide
                slide.shapes.add_picture(img_path, left, top, width, height)
            except FileNotFoundError:
                print(f"Image file not found: {img_path}")
            except Exception as e:
                print(f"Error adding image: {e}")
            image_index += 1

    # Save the presentation
    prs.save(output_filename)
    print(f"Presentation saved to {output_filename}")
    if os.path.exists('images'):
        try:
            for filename in os.listdir('images'):
                file_path = os.path.join('images', filename)
                os.remove(file_path)
            os.rmdir('images')
        except OSError as e:
            print(f"Failed to cleanup 'images' directory: {e}")

# Prompt user to select a template
print("Select a template:")
for key, path in TEMPLATES.items():
    print(f"{key}: {path}")

template_choice = int(input("Enter the number of the template you want to use: "))
template_path = TEMPLATES.get(template_choice, None)

# Example usage
topic = ['python']
create_presentation(gettext(topic, True), getphoto(gettext(topic, True)), template_path=template_path)
