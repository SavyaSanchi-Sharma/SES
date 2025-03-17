import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re
import json

def create_title_slide(prs, title):
    """Create a title slide with the main topic."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Set subtitle placeholder if available in the layout
    try:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = "Overview"
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(28)
        subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    except:
        pass  # Subtitle placeholder might not exist
    
    return slide

def create_content_slide(prs, title, content_bullets=None, code=None):
    """Create a content slide with optional bullet points and code."""
    slide_layout = prs.slide_layouts[1]  # Content slide layout with title and content
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Add bullet points if provided
    if content_bullets:
        try:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()  # Clear any default text
            
            for i, bullet in enumerate(content_bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = bullet
                
                # Determine bullet level based on the content
                if bullet.strip().startswith('*'):
                    p.level = 1
                    p.text = bullet.strip()[1:].strip()  # Remove the bullet marker
                else:
                    p.level = 0
                
                p.font.size = Pt(20)
        except:
            # If no placeholder, create a textbox
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(4.5)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            
            for i, bullet in enumerate(content_bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = bullet
                
                # Determine bullet level based on the content
                if bullet.strip().startswith('*'):
                    p.level = 1
                    p.text = bullet.strip()[1:].strip()  # Remove the bullet marker
                else:
                    p.level = 0
                
                p.font.size = Pt(20)
    
    # Add code if provided
    if code and isinstance(code, str) and code.strip():
        # Position code textbox
        left = Inches(1)
        top = Inches(5) if content_bullets else Inches(2)
        width = Inches(8)
        height = Inches(3)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = code.strip()
        p.font.size = Pt(14)
        p.font.name = "Courier New"
        
        # Add a light gray background to code box
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    return slide


def extract_title_from_point(point):
    """Extract a title from a numbered point."""
    # Try to match patterns like "1. Title: Content" or "1. Title - Content"
    title_match = re.match(r'^\d+\.\s+\*\*([^:]+)(?:\*\*)?[:|\-]', point)
    if title_match:
        return title_match.group(1).strip()
    
    # Try to match just the number and first few words
    simple_match = re.match(r'^\d+\.\s+\*\*([^*]+)\*\*', point)
    if simple_match:
        return simple_match.group(1).strip()
    
    # If no clear title format, just take the first part
    parts = re.match(r'^\d+\.\s+(.*?)(?:[:,]|$)', point)
    if parts:
        title = parts.group(1).strip()
        # Limit title length
        if len(title) > 50:
            return title[:47] + "..."
        return title
    
    # Fallback to numbered point
    num_match = re.match(r'^\d+\.', point)
    if num_match:
        return f"Point {num_match.group(0)}"
    
    return "Content Point"

def clean_text(text):
    """Clean text to remove markdown formatting."""
    # Remove markdown bold markers
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    return text

def process_data_and_create_presentation(data, output_filename='presentation.pptx'):
    """Process the data and create a presentation."""
    prs = Presentation()
    
    try:
        # Parse the input data if it's a string
        if isinstance(data, str):
            try:
                data = json.loads(data)
            except:
                # Try to evaluate as Python literal
                import ast
                try:
                    data = ast.literal_eval(data)
                except:
                    raise ValueError("Unable to parse input data")
        
        # Ensure we're working with a list of dictionaries
        if not isinstance(data, list):
            data = [data]
        
        for item in data:
            # Extract the topic, summary, and code
            topic = item.get('Topic', 'Presentation Topic')
            summary = item.get('Summary', [])
            code = item.get('Code', '')
            
            # Create title slide
            create_title_slide(prs, topic)
            
            # Process each point in the summary
            current_point_content = []
            current_point_title = ""
            
            i = 0
            while i < len(summary):
                point = summary[i]
                
                # Check if this is a new main point (starts with a number)
                if re.match(r'^\d+\.', point):
                    # If we have accumulated content, create a slide for it
                    if current_point_title and current_point_content:
                        create_content_slide(prs, current_point_title, current_point_content)
                    
                    # Extract the title from the new point
                    current_point_title = extract_title_from_point(point)
                    
                    # Clean the point text and start a new content list
                    current_point_content = [clean_text(point)]
                else:
                    # This is a continuation or sub-point of the current main point
                    current_point_content.append(clean_text(point))
                
                i += 1
            
            # Don't forget to create a slide for the last point
            if current_point_title and current_point_content:
                create_content_slide(prs, current_point_title, current_point_content)
            
            # Add code slide if provided
            if code.strip():
                create_content_slide(prs, "Code Example", None, code)
    
    except Exception as e:
        # Create an error slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Error Processing Data"
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"An error occurred while processing the data: {str(e)}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 0, 0)
    
    # Save the presentation
    prs.save(output_filename)
    print(f"Presentation saved as {output_filename}")
    return output_filename



